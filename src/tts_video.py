import os
import tempfile
import textwrap
from PIL import Image, ImageDraw, ImageFont

# Pillow compatibility
if not hasattr(Image, 'ANTIALIAS'):
    Image.ANTIALIAS = Image.Resampling.LANCZOS

from moviepy.editor import (
    ImageClip, AudioFileClip, concatenate_videoclips,
    concatenate_audioclips, CompositeAudioClip
)
import pyttsx3
from pptx import Presentation


# ---------------------------------------------
# Convert PPTX slides into image frames (title + summary + icon)
# ---------------------------------------------
def pptx_to_images(pptx_path, tmp_dir, icon_dir="assets/icons"):
    prs = Presentation(pptx_path)
    images = []

    # Load icons if available
    icons = []
    if os.path.isdir(icon_dir):
        icons = [os.path.join(icon_dir, f) for f in os.listdir(icon_dir)
                 if f.lower().endswith(('.png', '.jpg', '.jpeg'))]

    for i, slide in enumerate(prs.slides):
        img_path = os.path.join(tmp_dir, f"slide_{i+1}.png")

        # Create blank white image
        im = Image.new("RGB", (1280, 720), (255, 255, 255))
        draw = ImageDraw.Draw(im)

        # Load fonts (fallback-safe)
        try:
            title_font = ImageFont.truetype("arial.ttf", 50)
            note_font = ImageFont.truetype("arial.ttf", 30)
        except:
            title_font = ImageFont.load_default()
            note_font = ImageFont.load_default()

        # Extract first title and first text/note
        title, note = "", ""
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if not title:
                    title = shape.text.strip()
                else:
                    note = shape.text.strip()
                    break

        # Wrap text neatly
        wrapped_title = textwrap.fill(title, width=30)
        wrapped_note = textwrap.fill(note, width=60)

        # Draw title and note
        draw.text((60, 60), wrapped_title, fill=(0, 0, 0), font=title_font)
        draw.text((60, 200), wrapped_note, fill=(60, 60, 60), font=note_font)

        # Add icon (bottom-right)
        if icons:
            try:
                icon_path = icons[i % len(icons)]
                icon = Image.open(icon_path).convert("RGBA").resize((120, 120))
                im.paste(icon, (1120, 560), icon)
            except Exception as e:
                print(f"⚠️ Could not add icon for slide {i+1}: {e}")

        im.save(img_path)
        images.append(img_path)

    return images


# ---------------------------------------------
# Generate TTS audio for each slide
# ---------------------------------------------
def generate_tts_audio_per_slide(slides, tmp_dir):
    """
    Generates one audio file per slide using pyttsx3.
    Returns [{'path': <path>, 'duration': <float>}].
    """
    engine = pyttsx3.init()
    engine.setProperty("rate", 150)

    # Create per-slide audio
    for i, s in enumerate(slides):
        text = s.get("note", "").strip() or s.get("title", "")
        audio_path = os.path.join(tmp_dir, f"slide_{i+1}.mp3")
        engine.save_to_file(text, audio_path)
    engine.runAndWait()

    # Get real durations
    audio_info = []
    for i, s in enumerate(slides):
        audio_path = os.path.join(tmp_dir, f"slide_{i+1}.mp3")
        if os.path.exists(audio_path):
            clip = AudioFileClip(audio_path)
            duration = clip.duration
            clip.close()
            audio_info.append({"path": audio_path, "duration": duration})
        else:
            audio_info.append({"path": None, "duration": 0.0})
    return audio_info


# ---------------------------------------------
# Build final video with per-slide audio sync
# ---------------------------------------------
def make_video_from_pptx(pptx_path, slides, out_path, tts=True):
    tmp = tempfile.mkdtemp()
    images = pptx_to_images(pptx_path, tmp)

    if not tts:
        # fallback: simple no-TTS mode
        clips = [ImageClip(img).set_duration(5).resize(width=1280) for img in images]
        video = concatenate_videoclips(clips, method="compose")
        video.write_videofile(out_path, fps=24, codec="libx264")
        return

    try:
        # 1️⃣ Generate per-slide TTS audio and get durations
        audio_info = generate_tts_audio_per_slide(slides, tmp)

        # 2️⃣ Create one ImageClip per slide, attach its audio
        clips = []
        for i, img in enumerate(images):
            info = audio_info[i] if i < len(audio_info) else {"path": None, "duration": 4.0}
            dur = max(1.5, info["duration"])

            # fade = 12% of clip duration, clamped
            fade_time = max(0.15, min(1.5, dur * 0.12))
            fade_time = min(fade_time, dur / 2 - 0.01)

            clip = (
                ImageClip(img)
                .set_duration(dur)
                .resize(width=1280)
                .fadein(fade_time)
                .fadeout(fade_time)
            )

            # attach TTS audio if available
            if info["path"] and os.path.exists(info["path"]):
                audio_clip = AudioFileClip(info["path"])
                clip = clip.set_audio(audio_clip)

            clips.append(clip)

        # 3️⃣ Concatenate clips — keeps perfect per-slide audio alignment
        video = concatenate_videoclips(clips, method="compose")

        # 4️⃣ Add background music (soft)
        bg_path = os.path.join("assets", "background.mp3")
        if os.path.exists(bg_path):
            music = AudioFileClip(bg_path).volumex(0.05)
            total_dur = video.duration
            # loop music if needed
            if music.duration < total_dur:
                loops = int(total_dur // music.duration) + 1
                music = concatenate_audioclips([music] * loops)
            music = music.subclip(0, total_dur)
            final_audio = CompositeAudioClip([video.audio, music])
            video = video.set_audio(final_audio)

        # 5️⃣ Export video
        video.write_videofile(out_path, fps=24, codec="libx264", audio_codec="aac")

    except Exception as e:
        print("⚠️ Error during video creation:", e)
        fallback = [ImageClip(img).set_duration(5) for img in images]
        concatenate_videoclips(fallback, method="compose").write_videofile(out_path, fps=24, codec="libx264")
