from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os
from PIL import Image

def build_pptx(slides, out_path, template=None, icon_dir='assets/icons'):
    prs = Presentation()

    for s in slides:
        # use a blank layout
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # fully blank layout

        # --- Title (Topic Name) ---
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(8.5), Inches(1.5))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True  # ✅ enables wrapping
        title_p = title_tf.paragraphs[0]
        title_p.text = s['title']
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.alignment = PP_ALIGN.LEFT

        # --- One-line Summary / Note ---
        note_text = s.get('note', '')
        note_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.9), Inches(8.5), Inches(2.5))
        note_tf = note_box.text_frame
        note_tf.word_wrap = True  # ✅ enables wrapping
        note_p = note_tf.paragraphs[0]
        note_p.text = note_text
        note_p.font.size = Pt(20)
        note_p.font.bold = False
        note_p.alignment = PP_ALIGN.LEFT


        # --- Optional Icon / Visual ---
        image_path = find_icon_for_slide(s, icon_dir)
        if image_path:
            try:
                slide.shapes.add_picture(image_path, Inches(7.2), Inches(3.0), width=Inches(2.0))
            except Exception as e:
                print(f"Warning: could not add image for '{s['title']}': {e}")

    prs.save(out_path)


def find_icon_for_slide(slide_obj, icon_dir):
    """
    Picks an icon programmatically (or randomly) from assets/icons/
    """
    if not os.path.isdir(icon_dir):
        return None
    files = [os.path.join(icon_dir, f) for f in os.listdir(icon_dir)
             if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
    if not files:
        return None
    return files[hash(slide_obj['title']) % len(files)]
